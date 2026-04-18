#!/usr/bin/env python3
"""
ubcs-data-parser.py — Zoe's script
Đọc file Excel UBCS, xuất JSON chuẩn hóa vào /tmp/ubcs_data.json

Usage:
    python3 ubcs-data-parser.py <path_to_xlsx> [q4_template_pptx] [qd_phap_luat_xlsx]
"""
import sys
import json
import openpyxl
from collections import Counter, defaultdict
from pptx import Presentation

TRANG_THAI_MAP = {
    "giữ nguyên":           "Giữ nguyên",
    "đã điều chỉnh":        "Đã điều chỉnh",
    "điều chỉnh":           "SĐBS/Điều chỉnh",
    "đã thay thế":          "Đã thay thế",
    "thay thế":             "Thay thế",
    "sẽ tuyên hủy":         "Sẽ tuyên hủy",
    "tuyên hủy":            "Tuyên hủy",
    "hết hiệu lực":         "Hết hiệu lực",
    "xem xét lại":          "Xem xét lại",
    "đây là vbsp":          "VBSP (loại)",
    "bị lặp":               "Bị lặp",
    "vb khác":              "VB khác",
}

GROUP = {
    "Giữ nguyên":       "Giữ nguyên",
    "Đã điều chỉnh":    "SĐBS/Thay thế",
    "SĐBS/Điều chỉnh":  "SĐBS/Thay thế",
    "Đã thay thế":      "SĐBS/Thay thế",
    "Thay thế":         "SĐBS/Thay thế",
    "Tuyên hủy":        "Tuyên hủy",
    "Sẽ tuyên hủy":     "Tuyên hủy",
    "Hết hiệu lực":     "Tuyên hủy",
    "Xem xét lại":      "SĐBS/Thay thế",
    "VBSP (loại)":      "Tuyên hủy",
    "Bị lặp":           "Tuyên hủy",
    "VB khác":          "Giữ nguyên",
    "(chưa PL)":        "SĐBS/Thay thế",
}

# Mapping tên khối trong template Q4 → tên trong Excel
KHOI_MAP = {
    "Khối QLRR":        ["K QLRR TD", "K QLRR"],
    "Khối nhân sự":     ["K NHÂN SỰ"],
    "Khối tài chính":   ["k tài chính", "K TÀI CHÍNH"],
    "Khối PCTT":        ["k PCTT", "K PCTT"],
    "Khối Bán lẻ":      ["k Bán lẻ", "K BÁN LẺ"],
    "Khối KDV&TT":      ["k KDV &TT", "K KDV&TT"],
    "Văn phòng HĐQT":   ["VP HĐQT"],
    "Khối Vận hành":    ["K VẬN HÀNH"],
    "Khối KHDN":        ["K KHDN ", "K KHDN"],
    "Khối CNTT":        ["K CNTT"],
    "Khối DL&AI ":      ["k dữ liệu", "K DỮ LIỆU", "k dữ liệu"],
    "Khối mua sắm":     ["K Mua sắm", "K MUA SẮM"],
    "Các phòng còn lại (VP NHCT, KT nội bộ, P Pháp chế)": ["P PHÁP CHẾ", "P QLTT", "TT QLXLN"],
}

def norm_ts(v):
    if not v: return "(chưa PL)"
    s = str(v).strip().lower()
    for k, vv in TRANG_THAI_MAP.items():
        if k in s: return vv
    return str(v).strip()

def parse_phan1(wb, q4_khoi_map):
    """Đếm VBCS còn hiệu lực theo Khối từ sheet 'Gốc tổng hợp VPCS các khối'"""
    ws = wb["Gốc tổng hợp VPCS các khối"]
    per_khoi = Counter()
    for r in range(3, ws.max_row + 1):
        k   = ws.cell(r, 1).value
        cay = ws.cell(r, 6).value
        hl  = ws.cell(r, 9).value
        if not k: continue
        if cay != "Văn bản chính sách": continue
        if hl and "hết" in str(hl).lower(): continue
        per_khoi[str(k).strip()] += 1

    # Build per-khoi list với delta
    khoi_list = []
    for k, q1 in sorted(per_khoi.items(), key=lambda x: -x[1]):
        q4 = q4_khoi_map.get(k)
        delta = (q1 - q4) if q4 is not None else None
        khoi_list.append({
            "ten": k,
            "q1": q1,
            "q4": q4,
            "delta": delta,
        })

    return {
        "tong_q1": sum(per_khoi.values()),
        "theo_khoi": khoi_list,
    }

def parse_phan2(wb):
    """Phân loại trạng thái VB trên 5 năm theo Khối"""
    ws = wb["VB TRÊN 5N"]
    gn  = Counter()
    sd  = Counter()
    th  = Counter()
    tot = Counter()
    co_tien_do = []

    for r in range(3, ws.max_row + 1):
        k = ws.cell(r, 1).value
        if not k: continue
        trich = ws.cell(r, 3).value
        soky  = ws.cell(r, 4).value
        ts    = norm_ts(ws.cell(r, 15).value)
        tien  = ws.cell(r, 18).value
        ghi   = ws.cell(r, 19).value
        kk    = str(k).strip()
        grp   = GROUP.get(ts, "SĐBS/Thay thế")
        tot[kk] += 1
        if grp == "Giữ nguyên":      gn[kk] += 1
        elif grp == "SĐBS/Thay thế": sd[kk] += 1
        else:                         th[kk] += 1

        if tien:
            co_tien_do.append({
                "khoi":       kk,
                "soky":       str(soky or ""),
                "trich":      str(trich or ""),
                "trang_thai": ts,
                "tien_do":    str(tien),
                "ghi_chu":    str(ghi) if ghi else "",
            })

    khoi_list = [
        {
            "ten":        k,
            "giu_nguyen": gn[k],
            "sdbs":       sd[k],
            "tuyen_huy":  th[k],
            "tong":       tot[k],
        }
        for k, _ in sorted(tot.items(), key=lambda x: -x[1])
    ]
    co_tien_do.sort(key=lambda x: (x["khoi"], x["trang_thai"]))

    return {
        "tong": sum(tot.values()),
        "theo_khoi": khoi_list,
        "co_tien_do": co_tien_do,
    }

def extract_q4_from_template(template_path):
    """Lấy số VB Q4/2025 theo Khối từ slide 4 của template"""
    try:
        prs = Presentation(template_path)
        slide4 = prs.slides[3]
        q4_map = {}
        for shape in slide4.shapes:
            if shape.has_table:
                tbl = shape.table
                for r in range(1, len(tbl.rows)):
                    khoi = tbl.cell(r, 0).text.strip()
                    val  = tbl.cell(r, 3).text.strip()
                    if khoi.lower() in ("cộng", "tổng", ""): continue
                    try: q4_map[khoi] = int(val)
                    except: pass
        # Build reverse map: excel tên → q4 count
        excel_to_q4 = {}
        for template_name, excel_names in KHOI_MAP.items():
            q4_val = q4_map.get(template_name)
            if q4_val is None: continue
            for en in excel_names:
                excel_to_q4[en.strip().lower()] = q4_val
        return excel_to_q4
    except Exception as e:
        print(f"  [warn] Không đọc được Q4 từ template: {e}")
        return {}

def parse_phan3(qdpl_xlsx_path):
    """Đọc file Excel QĐPL, xuất danh sách VB pháp luật điều chỉnh VBCS."""
    wb = openpyxl.load_workbook(qdpl_xlsx_path, data_only=True)
    ws = wb.active
    rows = []
    tt_counter = 0
    for r in range(4, ws.max_row + 1):
        so   = ws.cell(r, 2).value
        ten  = ws.cell(r, 3).value
        if not so and not ten:
            continue
        tt_counter += 1
        tt        = ws.cell(r, 1).value or tt_counter
        ngay_bh   = ws.cell(r, 4).value
        ngay_hl   = ws.cell(r, 5).value
        noi_dung  = ws.cell(r, 6).value
        don_vi    = ws.cell(r, 7).value
        vb_ct     = ws.cell(r, 9).value
        ke_hoach  = ws.cell(r, 10).value
        note      = ws.cell(r, 11).value

        def fmt_date(v):
            if not v: return ""
            if hasattr(v, "strftime"): return v.strftime("%d/%m/%Y")
            s = str(v)
            if "00:00:00" in s: return s.replace(" 00:00:00", "").strip()
            return s.strip()

        rows.append({
            "tt":        str(tt),
            "so_vb":     str(so or "").strip(),
            "ten_vb":    str(ten or "").strip(),
            "ngay_bh":   fmt_date(ngay_bh),
            "ngay_hl":   fmt_date(ngay_hl),
            "noi_dung":  str(noi_dung or "").strip(),
            "don_vi":    str(don_vi or "").strip(),
            "vb_ct":     str(vb_ct or "").strip(),
            "ke_hoach":  str(ke_hoach or "").strip(),
            "note":      str(note or "").strip(),
        })
    return {"tong": len(rows), "danh_sach": rows}


def main():
    if len(sys.argv) < 2:
        print("Usage: python3 ubcs-data-parser.py <xlsx> [template.pptx] [qdpl.xlsx]")
        sys.exit(1)

    xlsx_path     = sys.argv[1]
    template_path = sys.argv[2] if len(sys.argv) > 2 else None
    qdpl_path     = sys.argv[3] if len(sys.argv) > 3 else None
    output_path   = "/tmp/ubcs_data.json"

    print(f"Đọc Excel: {xlsx_path}")
    wb = openpyxl.load_workbook(xlsx_path, data_only=True)

    # Extract Q4 reference từ template
    q4_khoi_map = {}
    if template_path:
        print(f"Đọc Q4 reference từ: {template_path}")
        raw_q4 = extract_q4_from_template(template_path)
        # Normalize: map excel khoi name → q4 count
        ws = wb["Gốc tổng hợp VPCS các khối"]
        excel_khois = set()
        for r in range(3, ws.max_row + 1):
            k = ws.cell(r, 1).value
            if k: excel_khois.add(str(k).strip())
        for ek in excel_khois:
            v = raw_q4.get(ek.lower())
            if v: q4_khoi_map[ek] = v

    # Detect quý/năm từ tên file
    import re, os
    fname = os.path.basename(xlsx_path)
    m = re.search(r'QU[YỲ]?\s*(\d)', fname, re.IGNORECASE)
    quy = m.group(1) if m else "?"
    m2 = re.search(r'(\d{4})', fname)
    nam = m2.group(1) if m2 else "2026"

    print("Phân tích Phần 1 (cây VBCS)...")
    phan1 = parse_phan1(wb, q4_khoi_map)
    print(f"  → {phan1['tong_q1']} VB, {len(phan1['theo_khoi'])} khối")

    print("Phân tích Phần 2 (VB > 5 năm)...")
    phan2 = parse_phan2(wb)
    print(f"  → {phan2['tong']} VB, {len(phan2['co_tien_do'])} VB có tiến độ")

    phan3 = None
    if qdpl_path:
        print(f"Phân tích Phần 3 (QĐPL điều chỉnh): {qdpl_path}")
        phan3 = parse_phan3(qdpl_path)
        print(f"  → {phan3['tong']} văn bản pháp luật")

    data = {
        "quy": quy,
        "nam": nam,
        "xlsx": xlsx_path,
        "phan1": phan1,
        "phan2": phan2,
        "phan3": phan3,
    }

    with open(output_path, "w", encoding="utf-8") as f:
        json.dump(data, f, ensure_ascii=False, indent=2)
    print(f"\nXuất: {output_path}")

if __name__ == "__main__":
    main()
