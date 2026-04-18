#!/usr/bin/env python3
"""
ubcs-slide-builder-v2.py
Layout theo PDF Gemini — white background, clean modern design.
Usage: python3 ubcs-slide-builder-v2.py <data.json> [output.pptx]
"""
import sys, json, os, copy
from datetime import date
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn
from pptx.chart.data import ChartData
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION
from lxml import etree

# ── Màu sắc ──────────────────────────────────
NAVY    = RGBColor(0x22, 0x3A, 0x5E)
BLUE    = RGBColor(0x19, 0x76, 0xD2)
ORANGE  = RGBColor(0xE6, 0x7E, 0x22)
RED     = RGBColor(0xC0, 0x39, 0x2B)
GREEN   = RGBColor(0x27, 0xAE, 0x60)
GRAY    = RGBColor(0x55, 0x55, 0x55)
LGRAY   = RGBColor(0xF2, 0xF4, 0xF7)
WHITE   = RGBColor(0xFF, 0xFF, 0xFF)
LBLUE   = RGBColor(0xE3, 0xF2, 0xFD)
BORDER  = RGBColor(0xDD, 0xDD, 0xDD)

CHART_COLORS = [
    RGBColor(0x19,0x76,0xD2), RGBColor(0xE6,0x7E,0x22),
    RGBColor(0xC0,0x39,0x2B), RGBColor(0x27,0xAE,0x60),
    RGBColor(0x95,0xA5,0xA6), RGBColor(0x8E,0x44,0xAD),
    RGBColor(0x16,0xA0,0x85), RGBColor(0xF3,0x9C,0x12),
    RGBColor(0x2C,0x3E,0x50), RGBColor(0x2E,0x86,0xC1),
]

FONT = "Arial"

SW = int(Inches(13.33))   # 16:9 widescreen
SH = int(Inches(7.5))

def rgb_t(r): return (r.red, r.green, r.blue)

# ── Helpers ───────────────────────────────────
def new_prs():
    prs = Presentation()
    prs.slide_width  = SW
    prs.slide_height = SH
    return prs

def blank_slide(prs):
    layout = prs.slide_layouts[6]   # completely blank
    s = prs.slides.add_slide(layout)
    # White background
    bg = s.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = WHITE
    return s

def txb(slide, text, left, top, width, height, *,
        size=14, bold=False, color=NAVY, align=PP_ALIGN.LEFT,
        italic=False, font=FONT, word_wrap=True):
    tb = slide.shapes.add_textbox(
        int(Inches(left)), int(Inches(top)),
        int(Inches(width)), int(Inches(height)))
    tf = tb.text_frame
    tf.word_wrap = word_wrap
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = str(text)
    r.font.size   = Pt(size)
    r.font.bold   = bold
    r.font.italic = italic
    r.font.color.rgb = color
    r.font.name   = font
    return tb

def hline(slide, top, left=0.5, width=12.33, color=BLUE, thickness_pt=2):
    """Đường kẻ ngang mỏng dưới tiêu đề."""
    h = slide.shapes.add_shape(
        1,
        int(Inches(left)), int(Inches(top)),
        int(Inches(width)), int(Pt(thickness_pt)))
    h.fill.solid(); h.fill.fore_color.rgb = color
    h.line.fill.background()
    return h

def rect(slide, left, top, width, height, fill_color, line_color=None):
    s = slide.shapes.add_shape(
        1,
        int(Inches(left)), int(Inches(top)),
        int(Inches(width)), int(Inches(height)))
    s.fill.solid(); s.fill.fore_color.rgb = fill_color
    if line_color:
        s.line.color.rgb = line_color
        s.line.width = Pt(0.5)
    else:
        s.line.fill.background()
    return s

def slide_title(slide, text, top=0.35):
    txb(slide, text, 0.5, top, 12.3, 0.8,
        size=28, bold=True, color=NAVY, font=FONT)
    hline(slide, top + 0.82)

def hex_str(color):
    """RGBColor → 'RRGGBB' hex string."""
    return str(color)

def set_cell_bg(cell, color):
    tc = cell._tc
    tcPr = tc.get_or_add_tcPr()
    for old in tcPr.findall(qn('a:solidFill')): tcPr.remove(old)
    fill = etree.SubElement(tcPr, qn('a:solidFill'))
    clr  = etree.SubElement(fill, qn('a:srgbClr'))
    clr.set('val', hex_str(color))

def set_cell_border(cell, color=BORDER):
    tc = cell._tc
    tcPr = tc.get_or_add_tcPr()
    ns = "http://schemas.openxmlformats.org/drawingml/2006/main"
    for side in ['lnL','lnR','lnT','lnB']:
        ln = tcPr.find(f"{{{ns}}}{side}")
        if ln is None:
            ln = etree.SubElement(tcPr, f"{{{ns}}}{side}")
        ln.set('w', '6350')  # 0.5pt
        solidFill = etree.SubElement(ln, f"{{{ns}}}solidFill")
        srgb = etree.SubElement(solidFill, f"{{{ns}}}srgbClr")
        srgb.set('val', hex_str(color))

def fill_cell(cell, text, *, bold=False, size=11, color=NAVY,
              align=PP_ALIGN.LEFT, bg=None, italic=False):
    cell.text = ""
    tf = cell.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.alignment = align
    r = p.add_run()
    r.text = str(text) if text is not None else ""
    r.font.size = Pt(size); r.font.bold = bold
    r.font.italic = italic
    r.font.color.rgb = color; r.font.name = FONT
    if bg: set_cell_bg(cell, bg)
    cell.margin_left  = Inches(0.08); cell.margin_right = Inches(0.08)
    cell.margin_top   = Inches(0.04); cell.margin_bottom = Inches(0.04)

def make_table(slide, nrows, ncols, left, top, width, height, col_widths=None):
    tbl = slide.shapes.add_table(
        nrows, ncols,
        int(Inches(left)), int(Inches(top)),
        int(Inches(width)), int(Inches(height))).table
    if col_widths:
        total = sum(col_widths)
        for i, w in enumerate(col_widths):
            tbl.columns[i].width = int(Inches(width) * w / total)
    return tbl

def style_header_row(tbl, bg=NAVY, fg=WHITE):
    for c in range(len(tbl.columns)):
        cell = tbl.cell(0, c)
        set_cell_bg(cell, bg)
        tf = cell.text_frame
        for para in tf.paragraphs:
            para.alignment = PP_ALIGN.CENTER
            for run in para.runs:
                run.font.bold = True; run.font.color.rgb = fg
                run.font.name = FONT; run.font.size = Pt(11)

def style_data_rows(tbl, start=1):
    for r in range(start, len(tbl.rows)):
        bg = LGRAY if r % 2 == 0 else WHITE
        for c in range(len(tbl.columns)):
            set_cell_bg(tbl.cell(r, c), bg)

def add_rows(tbl, n_data):
    """Đảm bảo bảng có đúng n_data data rows (không tính header)."""
    tbl_xml = tbl._tbl
    src_tr = tbl_xml.findall(qn('a:tr'))[1] if len(tbl.rows)>1 else tbl_xml.findall(qn('a:tr'))[0]
    while len(tbl.rows) < n_data + 1:
        tbl_xml.append(copy.deepcopy(src_tr))
    while len(tbl.rows) > n_data + 1:
        rows = tbl_xml.findall(qn('a:tr'))
        tbl_xml.remove(rows[-1])

def delta_color(v):
    try: v = float(str(v).replace('+',''))
    except: return GRAY
    if v > 0: return GREEN
    if v < 0: return RED
    return GRAY

# ── Slide 1: Title ────────────────────────────
def slide_01_title(prs, data):
    s = blank_slide(prs)
    quy, nam = data['quy'], data['nam']

    # Dải xanh góc phải trang trí
    deco = s.shapes.add_shape(1,
        int(SW - Inches(3)), 0, int(Inches(3)), int(SH))
    deco.fill.solid(); deco.fill.fore_color.rgb = RGBColor(0xF0,0xF7,0xFF)
    deco.line.fill.background()

    txb(s, "Báo cáo Cập nhật", 0.7, 2.2, 9, 0.9,
        size=36, bold=True, color=NAVY)
    txb(s, "Cây Văn bản Chính sách", 0.7, 3.1, 9, 1.0,
        size=42, bold=True, color=BLUE)
    txb(s, f"Kỳ báo cáo: Quý {quy}/{nam}  |  Phân tích Biến động & Tình trạng Văn bản trên 5 năm",
        0.7, 4.25, 10, 0.5, size=14, color=GRAY)
    hline(s, 4.15, left=0.7, width=5, color=BLUE, thickness_pt=3)

# ── Slide 2: Biến động Q1 vs Q4 ───────────────
def slide_02_biendong(prs, data):
    s = blank_slide(prs)
    slide_title(s, "BIẾN ĐỘNG SỐ LƯỢNG VB (Q1 VS Q4)")

    khois = data['phan1']['theo_khoi']
    # Lọc chỉ khối có Q4 reference
    rows_data = [(r['ten'], r['q1'], r['q4'], r['delta'])
                 for r in khois if r['q4'] is not None]
    # Thêm khối chỉ có Q1 (không có Q4)
    rows_no_q4 = [(r['ten'], r['q1'], None, None)
                  for r in khois if r['q4'] is None]

    all_rows = rows_data + rows_no_q4
    n = len(all_rows)

    tbl = make_table(s, n+1, 5, 0.5, 1.5, 12.3, 0.45 + 0.42*n,
                     col_widths=[3.5, 1.8, 1.8, 1.5, 3.7])
    headers = ["KHỐI/PHÒNG BAN", f"SL Q{data['quy']}/{data['nam']}",
               "SL Q4/2025", "TĂNG/GIẢM", "GHI CHÚ NGUYÊN NHÂN"]
    for c, h in enumerate(headers):
        fill_cell(tbl.cell(0,c), h, bold=True, size=10,
                  color=WHITE, bg=NAVY, align=PP_ALIGN.CENTER)

    for i, (ten, q1, q4, delta) in enumerate(all_rows):
        ri = i + 1
        bg = LGRAY if ri % 2 == 0 else WHITE
        fill_cell(tbl.cell(ri,0), ten, size=11, bg=bg)
        fill_cell(tbl.cell(ri,1), q1 if q1 else "—", size=11,
                  bold=True, align=PP_ALIGN.CENTER, bg=bg)
        fill_cell(tbl.cell(ri,2), q4 if q4 else "—", size=11,
                  align=PP_ALIGN.CENTER, bg=bg, color=GRAY)
        d_str = (f"+{delta}" if delta and delta>0 else str(delta)) if delta is not None else "—"
        dc = delta_color(delta) if delta else GRAY
        fill_cell(tbl.cell(ri,3), d_str, size=11, bold=True,
                  align=PP_ALIGN.CENTER, bg=bg, color=dc)
        fill_cell(tbl.cell(ri,4), "", size=10, bg=bg, color=GRAY, italic=True)

# ── Slide 3: Top 5 Khối bar chart ─────────────
def slide_03_top5(prs, data):
    s = blank_slide(prs)
    slide_title(s, "TOP 5 KHỐI CÓ KHỐI LƯỢNG VBCS LỚN NHẤT")

    khois = sorted(data['phan1']['theo_khoi'], key=lambda x: -x['q1'])[:5]
    tong  = data['phan1']['tong_q1']

    cd = ChartData()
    cd.categories = [r['ten'] for r in reversed(khois)]
    cd.add_series("Q1/2026", [r['q1'] for r in reversed(khois)])

    cf = s.shapes.add_chart(
        XL_CHART_TYPE.BAR_CLUSTERED,
        int(Inches(0.5)), int(Inches(1.5)),
        int(Inches(12.3)), int(Inches(5.5)),
        cd)
    chart = cf.chart
    chart.has_legend = False
    plot = chart.plots[0]
    plot.gap_width = 60

    for point in plot.series[0].points:
        point.format.fill.solid()
        point.format.fill.fore_color.rgb = BLUE

    chart.value_axis.has_major_gridlines = True
    chart.value_axis.major_gridlines.format.line.color.rgb = RGBColor(0xEE,0xEE,0xEE)
    chart.category_axis.tick_labels.font.size = Pt(12)
    chart.category_axis.tick_labels.font.color.rgb = NAVY

    # Note tổng
    txb(s, f"* Tổng toàn hệ thống có {tong} VBCS tại Quý {data['quy']}/{data['nam']}.",
        0.5, 6.9, 12, 0.4, size=11, color=GRAY, italic=True)

# ── Slide 4: Donut VB>5N ─────────────────────
def slide_04_donut(prs, data):
    s = blank_slide(prs)
    slide_title(s, "TỔNG QUAN VB TRÊN 5 NĂM (TRẠNG THÁI)")

    khois = data['phan2']['theo_khoi']
    gn   = sum(r['giu_nguyen'] for r in khois)
    sd   = sum(r['sdbs'] for r in khois)
    th   = sum(r['tuyen_huy'] for r in khois)
    tong = data['phan2']['tong']
    other = tong - gn - sd - th

    categories = ["Giữ nguyên", "Có kế hoạch điều chỉnh",
                  "Hết hiệu lực / Sẽ tuyên hủy", "Khác / Xem xét lại"]
    values     = [gn, sd, th, max(0, other)]

    cd = ChartData()
    cd.categories = categories
    cd.add_series("VB", values)

    cf = s.shapes.add_chart(
        XL_CHART_TYPE.DOUGHNUT,
        int(Inches(0.5)), int(Inches(1.4)),
        int(Inches(7.0)), int(Inches(5.8)),
        cd)
    chart = cf.chart
    chart.has_legend = True
    chart.legend.position = XL_LEGEND_POSITION.RIGHT
    chart.legend.include_in_layout = False

    donut_colors = [BLUE, ORANGE, RED, RGBColor(0x95,0xA5,0xA6)]
    for i, point in enumerate(chart.plots[0].series[0].points):
        point.format.fill.solid()
        point.format.fill.fore_color.rgb = donut_colors[i % len(donut_colors)]

    # Số tổng ở giữa donut
    txb(s, f"{tong} VB", 2.8, 3.8, 2.0, 0.7,
        size=22, bold=True, color=NAVY, align=PP_ALIGN.CENTER)

    # Legend thủ công bên phải với %
    legend_items = list(zip(categories, values, donut_colors))
    for i, (cat, val, col) in enumerate(legend_items):
        pct = round(val/tong*100) if tong else 0
        ly = 2.2 + i * 0.9
        r2 = rect(s, 7.8, ly+0.08, 0.25, 0.25, col)
        txb(s, f"{cat} ({val} VB - {pct}%)",
            8.2, ly, 4.5, 0.4, size=12, color=NAVY)

# ── Slide 5: Chi tiết trạng thái theo Khối ────
def slide_05_chitiet(prs, data):
    s = blank_slide(prs)
    slide_title(s, "CHI TIẾT TRẠNG THÁI VB > 5 NĂM THEO KHỐI")

    khois = sorted(data['phan2']['theo_khoi'], key=lambda x: -x['tong'])
    n = len(khois)

    tbl = make_table(s, n+2, 6, 0.5, 1.5, 12.3, 0.45 + 0.4*n,
                     col_widths=[3.2, 1.6, 1.8, 2.2, 1.8, 1.8])
    headers = ["KHỐI BAN HÀNH","GIỮ NGUYÊN","ĐIỀU CHỈNH","TUYÊN HỦY/HẾT HL","THAY THẾ","TỔNG CỘNG"]
    for c, h in enumerate(headers):
        fill_cell(tbl.cell(0,c), h, bold=True, size=10,
                  color=WHITE, bg=NAVY, align=PP_ALIGN.CENTER)

    gn_tot=sd_tot=th_tot=tt_tot=0
    for i, row in enumerate(khois):
        ri = i+1
        bg = LGRAY if ri%2==0 else WHITE
        fill_cell(tbl.cell(ri,0), row['ten'], size=11, bg=bg)
        fill_cell(tbl.cell(ri,1), row['giu_nguyen'] or "—", size=11,
                  align=PP_ALIGN.CENTER, bg=bg)
        # Điều chỉnh → orange bold
        fill_cell(tbl.cell(ri,2), row['sdbs'] or "—", size=11, bold=True,
                  align=PP_ALIGN.CENTER, bg=bg,
                  color=ORANGE if row['sdbs'] else GRAY)
        fill_cell(tbl.cell(ri,3), row['tuyen_huy'] or "—", size=11,
                  align=PP_ALIGN.CENTER, bg=bg,
                  color=RED if row['tuyen_huy'] else GRAY)
        fill_cell(tbl.cell(ri,4), "—", size=11, align=PP_ALIGN.CENTER, bg=bg, color=GRAY)
        fill_cell(tbl.cell(ri,5), row['tong'], size=11, bold=True,
                  align=PP_ALIGN.CENTER, bg=bg)
        gn_tot+=row['giu_nguyen']; sd_tot+=row['sdbs']
        th_tot+=row['tuyen_huy']; tt_tot+=row['tong']

    # Tổng
    last = n+1
    for c in range(6): set_cell_bg(tbl.cell(last,c), NAVY)
    fill_cell(tbl.cell(last,0),"TỔNG CỘNG",bold=True,size=11,color=WHITE,bg=NAVY)
    fill_cell(tbl.cell(last,1),gn_tot,bold=True,size=11,color=WHITE,bg=NAVY,align=PP_ALIGN.CENTER)
    fill_cell(tbl.cell(last,2),sd_tot,bold=True,size=11,color=WHITE,bg=NAVY,align=PP_ALIGN.CENTER)
    fill_cell(tbl.cell(last,3),th_tot,bold=True,size=11,color=WHITE,bg=NAVY,align=PP_ALIGN.CENTER)
    fill_cell(tbl.cell(last,4),"—",bold=True,size=11,color=WHITE,bg=NAVY,align=PP_ALIGN.CENTER)
    fill_cell(tbl.cell(last,5),tt_tot,bold=True,size=11,color=WHITE,bg=NAVY,align=PP_ALIGN.CENTER)

# ── Slide 6: Cards kế hoạch trọng điểm ────────
def slide_06_cards(prs, data):
    s = blank_slide(prs)
    slide_title(s, "KẾ HOẠCH ĐIỀU CHỈNH VĂN BẢN TRỌNG ĐIỂM")

    rows = [r for r in data['phan2']['co_tien_do'] if r['tien_do'].strip()]
    # Lấy top 3 VB có tiến độ rõ ràng nhất (ưu tiên có kế hoạch cụ thể)
    top3 = rows[:3]

    card_w, card_h = 3.8, 4.5
    card_tops = 1.6
    positions = [0.4, 4.5, 8.6]
    card_color = RGBColor(0xF8,0xFB,0xFF)
    accent_colors = [BLUE, ORANGE, RGBColor(0x27,0xAE,0x60)]

    for i, (left, row) in enumerate(zip(positions, top3)):
        acc = accent_colors[i % len(accent_colors)]
        # Card background
        card = rect(s, left, card_tops, card_w, card_h,
                    card_color, BORDER)
        # Accent bar trên
        rect(s, left, card_tops, card_w, 0.06, acc)
        # Khối
        txb(s, row['khoi'], left+0.15, card_tops+0.18, card_w-0.3, 0.35,
            size=13, bold=True, color=acc)
        # Tên VB (rút ngắn)
        ten = row['trich'][:80] + ("…" if len(row['trich'])>80 else "")
        txb(s, ten, left+0.15, card_tops+0.55, card_w-0.3, 0.8,
            size=10, color=NAVY)
        # Tiến độ badge
        tien = row['tien_do'][:60]
        badge = rect(s, left+0.15, card_tops+1.45, card_w-0.3, 0.32,
                     RGBColor(0xE8,0xF4,0xFF), acc)
        txb(s, f"Tiến độ: {tien}", left+0.2, card_tops+1.5, card_w-0.4, 0.28,
            size=9, color=BLUE, bold=True)
        # Ghi chú
        ghi = row['ghi_chu'][:120] if row['ghi_chu'] else ""
        txb(s, ghi, left+0.15, card_tops+1.9, card_w-0.3, 1.6,
            size=9, color=GRAY, italic=True)

# ── Slide 7: Tiến độ xử lý (bảng chi tiết) ───
def slide_07_tiendo(prs, data, page=0, per_page=8):
    rows = [r for r in data['phan2']['co_tien_do'] if r['tien_do'].strip()]
    chunk = rows[page*per_page:(page+1)*per_page]
    total_pages = max(1, (len(rows)+per_page-1)//per_page)

    s = blank_slide(prs)
    title_suffix = f" ({page+1}/{total_pages})" if total_pages > 1 else ""
    slide_title(s, f"TIẾN ĐỘ XỬ LÝ CỤ THỂ (TRÍCH XUẤT DỮ LIỆU EXCEL){title_suffix}")

    n = len(chunk)
    if not n: return s

    tbl = make_table(s, n+1, 4, 0.5, 1.5, 12.3, 0.4 + 0.48*n,
                     col_widths=[4.0, 2.0, 1.8, 4.5])
    headers = ["TÊN / SỐ KÝ HIỆU VĂN BẢN", "KHỐI", "TRẠNG THÁI", "GHI CHÚ / KẾ HOẠCH CHI TIẾT"]
    for c, h in enumerate(headers):
        fill_cell(tbl.cell(0,c), h, bold=True, size=10,
                  color=WHITE, bg=NAVY, align=PP_ALIGN.CENTER)

    STATUS_COLOR = {
        "Hết hiệu lực": RED, "Tuyên hủy": RED, "Sẽ tuyên hủy": RED,
        "SĐBS/Điều chỉnh": ORANGE, "Đã điều chỉnh": GREEN,
        "Thay thế": ORANGE, "Đã thay thế": GREEN,
        "Xem xét lại": RGBColor(0xE6,0x7E,0x22),
    }
    for i, row in enumerate(chunk):
        ri = i+1
        bg = LGRAY if ri%2==0 else WHITE
        label = f"{row['soky']}\n({row['trich'][:50]}…)" if len(row['trich'])>50 else f"{row['soky']}\n{row['trich']}"
        fill_cell(tbl.cell(ri,0), label, size=9, bg=bg, italic=True)
        fill_cell(tbl.cell(ri,1), row['khoi'], size=9, bg=bg)
        sc = STATUS_COLOR.get(row['trang_thai'], GRAY)
        fill_cell(tbl.cell(ri,2), row['trang_thai'], size=9, bold=True,
                  align=PP_ALIGN.CENTER, bg=bg, color=sc)
        ke = row['tien_do']
        if row['ghi_chu']: ke += "\n" + row['ghi_chu']
        fill_cell(tbl.cell(ri,3), ke[:200], size=9, bg=bg, color=GRAY)
    return s

# ── Slide 8: QĐPL Phần 3 ─────────────────────
def slide_08_qdpl(prs, data, page=0, per_page=6):
    if not data.get('phan3'): return None
    rows  = data['phan3']['danh_sach']
    chunk = rows[page*per_page:(page+1)*per_page]
    total = max(1,(len(rows)+per_page-1)//per_page)
    if not chunk: return None

    s = blank_slide(prs)
    suffix = f" ({page+1}/{total})" if total>1 else ""
    slide_title(s, f"VBCS ĐIỀU CHỈNH THEO QUY ĐỊNH PHÁP LUẬT QUÝ 1/2026{suffix}")

    n = len(chunk)
    tbl = make_table(s, n+1, 5, 0.5, 1.5, 12.3, 0.4+0.55*n,
                     col_widths=[0.6, 3.4, 2.4, 2.9, 3.0])
    headers = ["TT","TÊN VĂN BẢN PHÁP LUẬT","ĐƠN VỊ ĐẦU MỐI","VBCS CHỊU TÁC ĐỘNG","KẾ HOẠCH ĐIỀU CHỈNH"]
    for c, h in enumerate(headers):
        fill_cell(tbl.cell(0,c), h, bold=True, size=10,
                  color=WHITE, bg=NAVY, align=PP_ALIGN.CENTER)
    for i, row in enumerate(chunk):
        ri = i+1
        bg = LGRAY if ri%2==0 else WHITE
        ten = f"{row['so_vb']}\n{row['ten_vb'][:60]}" if row['so_vb'] else row['ten_vb'][:60]
        if row['ngay_hl']: ten += f"\n(HL: {row['ngay_hl']})"
        fill_cell(tbl.cell(ri,0), page*per_page+i+1, size=10, align=PP_ALIGN.CENTER, bg=bg)
        fill_cell(tbl.cell(ri,1), ten, size=9, bg=bg)
        fill_cell(tbl.cell(ri,2), row['don_vi'][:80], size=9, bg=bg, color=GRAY)
        fill_cell(tbl.cell(ri,3), row['vb_ct'][:100], size=9, bg=bg)
        kh = row['ke_hoach'] or row['note'] or "—"
        kh_color = BLUE if len(kh)>5 and kh!="—" else GRAY
        fill_cell(tbl.cell(ri,4), kh[:120], size=9, bg=bg, color=kh_color)
    return s

# ── Slide cuối: Tổng kết ─────────────────────
def slide_final(prs, data):
    s = blank_slide(prs)
    slide_title(s, "TỔNG KẾT NHIỆM VỤ TRỌNG TÂM")

    khois = data['phan2']['theo_khoi']
    sd_tot = sum(r['sdbs'] for r in khois)
    th_tot = sum(r['tuyen_huy'] for r in khois)
    tong5  = data['phan2']['tong']

    # Số lớn bên trái
    txb(s, str(sd_tot), 0.6, 2.4, 2.2, 1.4,
        size=72, bold=True, color=NAVY, align=PP_ALIGN.CENTER)
    txb(s, "Văn bản cần\nĐiều chỉnh", 0.6, 3.8, 2.2, 0.8,
        size=14, color=GRAY, align=PP_ALIGN.CENTER)

    hline(s, 2.3, left=3.2, width=0.05, color=BORDER, thickness_pt=80)  # divider dọc

    # Bullet points
    txb(s, "Trọng tâm phối hợp thực hiện", 3.6, 2.0, 9.2, 0.5,
        size=16, bold=True, color=NAVY)

    bullets = [
        f"Đôn đốc các phòng ban đầu mối (đặc biệt Khối QLRR, Khối Tài chính) bám sát deadline đã cam kết trong kế hoạch.",
        f"Thực hiện thao tác Tuyên hủy / Đánh dấu Hết hiệu lực trực tiếp trên hệ thống Intranet đối với {th_tot} văn bản không còn giá trị.",
        f"Tổ chức bàn giao rõ ràng đối với các văn bản 'Xem xét lại' để tìm đúng đơn vị nghiệp vụ thụ lý.",
    ]
    for i, b in enumerate(bullets):
        dot = rect(s, 3.6, 2.75 + i*0.85, 0.15, 0.15,
                   BLUE)
        txb(s, b, 3.9, 2.68 + i*0.85, 9.0, 0.7,
            size=12, color=NAVY)

# ── Main ──────────────────────────────────────
def main():
    if len(sys.argv) < 2:
        print("Usage: python3 ubcs-slide-builder-v2.py <data.json> [output.pptx]")
        sys.exit(1)

    data_path = sys.argv[1]
    with open(data_path, encoding='utf-8') as f:
        data = json.load(f)

    quy, nam = data['quy'], data['nam']
    today = date.today().strftime("%Y-%m-%d")
    out = sys.argv[2] if len(sys.argv)>2 else \
          os.path.expanduser(f"~/Downloads/Slide UBCS Quy {quy} {nam} v2 - {today}.pptx")

    prs = new_prs()
    n = [1]

    print(f"Build deck Quý {quy}/{nam}...")

    slide_01_title(prs, data)
    print(f"  ✓ Slide {n[0]}: title"); n[0]+=1

    slide_02_biendong(prs, data)
    print(f"  ✓ Slide {n[0]}: biến động Q1 vs Q4"); n[0]+=1

    slide_03_top5(prs, data)
    print(f"  ✓ Slide {n[0]}: top 5 khối bar chart"); n[0]+=1

    slide_04_donut(prs, data)
    print(f"  ✓ Slide {n[0]}: donut VB>5N"); n[0]+=1

    slide_05_chitiet(prs, data)
    print(f"  ✓ Slide {n[0]}: chi tiết trạng thái theo Khối"); n[0]+=1

    slide_06_cards(prs, data)
    print(f"  ✓ Slide {n[0]}: cards kế hoạch trọng điểm"); n[0]+=1

    # Tiến độ — có thể nhiều trang
    rows = [r for r in data['phan2']['co_tien_do'] if r['tien_do'].strip()]
    per_page = 8
    pages_td = max(1,(len(rows)+per_page-1)//per_page)
    for pi in range(pages_td):
        slide_07_tiendo(prs, data, page=pi, per_page=per_page)
        print(f"  ✓ Slide {n[0]}: tiến độ trang {pi+1}/{pages_td}"); n[0]+=1

    # Phần 3 — QĐPL
    if data.get('phan3'):
        rows3 = data['phan3']['danh_sach']
        pages3 = max(1,(len(rows3)+5)//6)
        for pi in range(pages3):
            slide_08_qdpl(prs, data, page=pi)
            print(f"  ✓ Slide {n[0]}: QĐPL trang {pi+1}/{pages3}"); n[0]+=1

    slide_final(prs, data)
    print(f"  ✓ Slide {n[0]}: tổng kết")

    prs.save(out)
    print(f"\n✓ Saved: {out}")
    print(f"  Total: {len(prs.slides)} slides")

if __name__ == '__main__':
    main()
