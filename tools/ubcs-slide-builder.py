#!/usr/bin/env python3
"""
ubcs-slide-builder.py — Katarina's script
Nhận JSON từ ubcs-data-parser, clone template PPTX, build slide UBCS đẹp.

Usage:
    python3 ubcs-slide-builder.py <data.json> <template.pptx> [output.pptx]
"""
import sys
import json
import copy
import os
from datetime import date
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.oxml.ns import qn
from pptx.enum.text import PP_ALIGN
from pptx.chart.data import ChartData
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION
from lxml import etree

STYLE_PATH = os.path.join(os.path.dirname(__file__), "ubcs-style-guide.json")
with open(STYLE_PATH, encoding="utf-8") as f:
    STYLE = json.load(f)

def rgb(hex_str):
    h = hex_str.lstrip("#")
    return RGBColor(int(h[0:2],16), int(h[2:4],16), int(h[4:6],16))

def rgb_tuple(hex_str):
    h = hex_str.lstrip("#")
    return (int(h[0:2],16), int(h[2:4],16), int(h[4:6],16))

C = {k: rgb(v) for k, v in STYLE["colors"].items()}
CT = {k: rgb_tuple(v) for k, v in STYLE["colors"].items()}

TITLE_FONT   = STYLE["fonts"]["title"]
BODY_FONT    = STYLE["fonts"]["body"]
FS           = STYLE["font_sizes"]
TS_COLOR_MAP = STYLE["trang_thai_color"]
CHART_COLORS = STYLE["chart_colors"]

# Style guide sections with fallback defaults
_header_bar    = STYLE.get("header_bar", {})
_table_style   = STYLE.get("table_style", {})
_slide_layouts = STYLE.get("slide_layouts", {})
_delta_style   = STYLE.get("delta", {})

HEADER_HEIGHT  = _header_bar.get("height_inches", 0.55)
ROWS_PER_PAGE  = _slide_layouts.get("detail_slide", {}).get("rows_per_page", 8)
DELTA_POS_COLOR = rgb(_delta_style.get("positive_color", STYLE["colors"]["green"]))
DELTA_NEG_COLOR = rgb(_delta_style.get("negative_color", STYLE["colors"]["red"]))

# ─────────────────────────────────────────────
# Core helpers
# ─────────────────────────────────────────────

def set_cell_bg(tc_el, rgb_tup):
    tcPr = tc_el.get_or_add_tcPr()
    for old in tcPr.findall(qn('a:solidFill')):
        tcPr.remove(old)
    fill = etree.SubElement(tcPr, qn('a:solidFill'))
    clr  = etree.SubElement(fill, qn('a:srgbClr'))
    clr.set('val', '{:02X}{:02X}{:02X}'.format(*rgb_tup))

def set_cell(cell, text, *, bold=False, size=None, bg=None, fg=None,
             align=None, font=None, italic=False):
    cell.text = ""
    tf = cell.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    if align: p.alignment = align
    run = p.add_run()
    run.text = str(text) if text is not None else ""
    run.font.size  = Pt(size or FS["table_body"])
    run.font.bold  = bold
    run.font.italic = italic
    run.font.name  = font or BODY_FONT
    if fg: run.font.color.rgb = fg
    if bg: set_cell_bg(cell._tc, rgb_tuple(bg) if isinstance(bg, str) else bg)
    cell.margin_left   = Inches(0.06)
    cell.margin_right  = Inches(0.06)
    cell.margin_top    = Inches(0.03)
    cell.margin_bottom = Inches(0.03)

def set_cell_text_keep_format(tc_el, new_text):
    """Thay text nhưng giữ nguyên rPr (format) của run đầu tiên."""
    ns = "http://schemas.openxmlformats.org/drawingml/2006/main"
    paras = tc_el.findall(f".//{{{ns}}}p")
    if not paras: return
    first_rPr = None
    for p in paras:
        for r in p.findall(f"{{{ns}}}r"):
            if first_rPr is None:
                rp = r.find(f"{{{ns}}}rPr")
                if rp is not None: first_rPr = copy.deepcopy(rp)
            p.remove(r)
    for p in paras[1:]:
        paras[0].getparent().remove(p)
    p0 = paras[0]
    r_el = etree.SubElement(p0, f"{{{ns}}}r")
    if first_rPr is not None: r_el.insert(0, first_rPr)
    t_el = etree.SubElement(r_el, f"{{{ns}}}t")
    t_el.text = str(new_text) if new_text is not None else ""

def clone_slide(src_prs, slide_idx, dest_prs):
    src_slide = src_prs.slides[slide_idx]
    blank     = dest_prs.slide_layouts[6]
    new_slide = dest_prs.slides.add_slide(blank)
    sp_tree   = new_slide.shapes._spTree
    for el in list(sp_tree): sp_tree.remove(el)
    for el in src_slide.shapes._spTree:
        sp_tree.append(copy.deepcopy(el))
    try:
        src_bg  = src_slide.background._element
        dest_bg = new_slide.background._element
        for child in list(dest_bg): dest_bg.remove(child)
        for child in src_bg: dest_bg.append(copy.deepcopy(child))
    except Exception: pass
    return new_slide

def get_table(slide):
    for s in slide.shapes:
        if s.has_table: return s.table, s
    return None, None

def add_table(slide, nrows, ncols, left, top, width, height, col_widths=None):
    tbl = slide.shapes.add_table(nrows, ncols,
                                  int(left), int(top),
                                  int(width), int(height)).table
    if col_widths:
        total = sum(col_widths)
        for i, w in enumerate(col_widths):
            tbl.columns[i].width = int(width * w / total)
    return tbl

def add_rows_to_table(tbl, n_need):
    """Thêm/bớt data rows (giữ header row 0) để đủ n_need data rows."""
    tbl_xml = tbl._tbl
    tr_list = tbl_xml.findall(qn('a:tr'))
    src_tr  = tr_list[1] if len(tr_list) > 1 else tr_list[0]
    # Xoá bớt
    while len(tbl.rows) > n_need + 1:
        rows = tbl_xml.findall(qn('a:tr'))
        tbl_xml.remove(rows[-1])
    # Thêm
    while len(tbl.rows) < n_need + 1:
        tbl_xml.append(copy.deepcopy(src_tr))

def delta_str(v):
    if v is None: return "—"
    if v > 0: return f"+{v}"
    if v < 0: return str(v)
    return "0"

def delta_color(v):
    if v is None: return None
    if v > 0: return DELTA_POS_COLOR
    if v < 0: return DELTA_NEG_COLOR
    return None

_hb_title  = _header_bar.get("title",       {})
_hb_secnum = _header_bar.get("section_num", {})

def add_header(slide, title, subtitle=None, section_num=None, SW=None):
    SW = SW or slide.part.package.presentation.slide_width
    bar_color = _header_bar.get("color", STYLE["colors"]["navy"])
    # Dải navy trên cùng
    bar = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        0, 0, int(SW), int(Inches(HEADER_HEIGHT))
    )
    bar.fill.solid()
    bar.fill.fore_color.rgb = rgb(bar_color)
    bar.line.fill.background()

    # Section number nhỏ góc trái
    if section_num:
        sn_left = _hb_secnum.get("left", 0.2)
        sn_top  = _hb_secnum.get("top",  0.1)
        sn_size = _hb_secnum.get("size", FS.get("section_label", 18))
        sn_font = _hb_secnum.get("font", TITLE_FONT)
        tb = slide.shapes.add_textbox(Inches(sn_left), Inches(sn_top), Inches(0.6), Inches(0.4))
        tf = tb.text_frame; p = tf.paragraphs[0]
        r = p.add_run(); r.text = str(section_num)
        r.font.size = Pt(sn_size); r.font.bold = _hb_secnum.get("bold", True)
        r.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
        r.font.name = sn_font

    # Title
    t_left = _hb_title.get("left", 0.7 if section_num else 0.2)
    t_top  = _hb_title.get("top",  0.06)
    t_size = _hb_title.get("size", FS["slide_title"])
    t_font = _hb_title.get("font", TITLE_FONT)
    tb = slide.shapes.add_textbox(
        Inches(t_left if section_num else 0.2), Inches(t_top),
        int(SW - Inches(1.2)), Inches(HEADER_HEIGHT - 0.1)
    )
    tf = tb.text_frame; p = tf.paragraphs[0]
    r = p.add_run(); r.text = title
    r.font.size = Pt(t_size); r.font.bold = _hb_title.get("bold", True)
    r.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
    r.font.name = t_font

    # Subtitle
    if subtitle:
        tb2 = slide.shapes.add_textbox(Inches(0.2), Inches(0.58), int(SW - Inches(0.4)), Inches(0.35))
        tf2 = tb2.text_frame; p2 = tf2.paragraphs[0]
        r2 = p2.add_run(); r2.text = subtitle
        r2.font.size = Pt(10); r2.font.italic = True
        r2.font.color.rgb = C["gray_text"]; r2.font.name = BODY_FONT

def style_table_header(tbl):
    """Apply navy header + alternating rows."""
    ncols = len(tbl.columns)
    for c in range(ncols):
        cell = tbl.cell(0, c)
        set_cell_bg(cell._tc, CT["navy"])
        tf = cell.text_frame
        for para in tf.paragraphs:
            para.alignment = PP_ALIGN.CENTER
            for run in para.runs:
                run.font.bold  = True
                run.font.color.rgb = C["white"]
                run.font.name  = TITLE_FONT
                run.font.size  = Pt(FS["table_header"])

def style_data_row(tbl, row_idx, bg=None):
    default_bg = CT["light_blue"] if row_idx % 2 == 0 else None
    final_bg = rgb_tuple(bg) if isinstance(bg, str) else (bg or default_bg)
    if final_bg:
        for c in range(len(tbl.columns)):
            set_cell_bg(tbl.cell(row_idx, c)._tc, final_bg)

def style_total_row(tbl, row_idx):
    for c in range(len(tbl.columns)):
        cell = tbl.cell(row_idx, c)
        set_cell_bg(cell._tc, CT["navy"])
        tf = cell.text_frame
        for para in tf.paragraphs:
            for run in para.runs:
                run.font.bold  = True
                run.font.color.rgb = C["white"]
                run.font.name  = TITLE_FONT
                run.font.size  = Pt(FS["table_header"])

# ─────────────────────────────────────────────
# Slide builders
# ─────────────────────────────────────────────

def build_title_slide(src_prs, dest_prs, data):
    s = clone_slide(src_prs, 0, dest_prs)
    quy, nam = data["quy"], data["nam"]
    for shape in s.shapes:
        if shape.has_text_frame:
            for para in shape.text_frame.paragraphs:
                for run in para.runs:
                    t = run.text
                    if "tháng" in t.lower() or "/202" in t or "quý" in t.lower():
                        run.text = f"Quý {quy}/{nam}"
    return s

def build_muc_luc(src_prs, dest_prs, data, SW, SH):
    """Slide danh mục 3 phần, clone từ slide 2 rồi thêm nội dung."""
    s = clone_slide(src_prs, 1, dest_prs)
    has_phan3 = bool(data.get("phan3"))

    muc = [
        ("1", "Cây Văn Bản Chính Sách",
         f"Tổng {data['phan1']['tong_q1']} VBCS còn hiệu lực theo {len(data['phan1']['theo_khoi'])} Khối"),
        ("2", "Văn Bản Chính Sách Trên 5 Năm",
         f"Tổng {data['phan2']['tong']} VB — Giữ nguyên / SĐBS / Tuyên hủy theo Khối"),
    ]
    if has_phan3:
        muc.append(("3", "VBCS Điều Chỉnh Theo QĐPL Quý 1/2026",
                    f"{data['phan3']['tong']} văn bản pháp luật có tác động đến VBCS"))

    # Vị trí các mục — căn giữa dọc
    total_h = len(muc) * 1.35
    start_top = (SH / 914400 - total_h) / 2 + 0.2

    for i, (num, title, subtitle) in enumerate(muc):
        top = start_top + i * 1.35

        # Số mục — hình tròn navy
        circle = s.shapes.add_shape(9, int(Inches(0.5)), int(Inches(top)),
                                     int(Inches(0.7)), int(Inches(0.7)))
        circle.fill.solid(); circle.fill.fore_color.rgb = C["navy"]
        circle.line.fill.background()
        tf = circle.text_frame; tf.margin_top = 0; tf.margin_bottom = 0
        p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
        r = p.add_run(); r.text = num
        r.font.size = Pt(22); r.font.bold = True
        r.font.color.rgb = C["white"]; r.font.name = TITLE_FONT

        # Tiêu đề mục
        tb = s.shapes.add_textbox(int(Inches(1.35)), int(Inches(top - 0.03)),
                                   int(SW - Inches(1.6)), int(Inches(0.45)))
        tf = tb.text_frame; p = tf.paragraphs[0]
        r = p.add_run(); r.text = title
        r.font.size = Pt(20); r.font.bold = True
        r.font.color.rgb = C["navy"]; r.font.name = TITLE_FONT

        # Subtitle mô tả
        tb2 = s.shapes.add_textbox(int(Inches(1.35)), int(Inches(top + 0.42)),
                                    int(SW - Inches(1.6)), int(Inches(0.4)))
        tf2 = tb2.text_frame; p2 = tf2.paragraphs[0]
        r2 = p2.add_run(); r2.text = subtitle
        r2.font.size = Pt(11); r2.font.color.rgb = C["gray_text"]
        r2.font.name = BODY_FONT; r2.font.italic = True

        # Đường kẻ phân cách (trừ mục cuối)
        if i < len(muc) - 1:
            line = s.shapes.add_shape(1,
                int(Inches(0.4)), int(Inches(top + 0.95)),
                int(SW - Inches(0.8)), int(Pt(1)))
            line.fill.solid(); line.fill.fore_color.rgb = C["light_blue"]
            line.line.fill.background()

    return s


def build_section_divider(src_prs, dest_prs, num, title, SW, SH):
    """Slide bìa phần — clone slide 2 template."""
    s = clone_slide(src_prs, 1, dest_prs)
    # Dải navy trái
    bar = s.shapes.add_shape(1, 0, 0, int(Inches(0.2)), int(SH))
    bar.fill.solid(); bar.fill.fore_color.rgb = C["navy"]
    bar.line.fill.background()

    tb = s.shapes.add_textbox(int(Inches(0.5)), int(SH/2 - Inches(0.8)),
                               int(SW - Inches(0.8)), int(Inches(1.6)))
    tf = tb.text_frame
    p1 = tf.paragraphs[0]; p1.alignment = PP_ALIGN.LEFT
    r1 = p1.add_run(); r1.text = f"PHẦN {num}"
    r1.font.size = Pt(28); r1.font.bold = True
    r1.font.color.rgb = C["blue"]; r1.font.name = TITLE_FONT

    p2 = tf.add_paragraph(); p2.alignment = PP_ALIGN.LEFT
    r2 = p2.add_run(); r2.text = title
    r2.font.size = Pt(36); r2.font.bold = True
    r2.font.color.rgb = C["navy"]; r2.font.name = TITLE_FONT
    return s

def build_phan1_bang(src_prs, dest_prs, data, SW):
    """Slide: bảng cây VBCS theo Khối"""
    s = clone_slide(src_prs, 3, dest_prs)  # clone slide 4
    tbl, _ = get_table(s)
    quy, nam = data["quy"], data["nam"]

    # Cập nhật header columns
    set_cell_text_keep_format(tbl.cell(0, 3)._tc, f"Q{quy}/{nam}")
    set_cell_text_keep_format(tbl.cell(0, 4)._tc, f"So sánh Q{quy}/{nam}\nvới 31/12/{int(nam)-1}")

    tong_q4 = 0
    tong_q1 = 0
    khoi_dict = {row["ten"].strip().lower(): row for row in data["phan1"]["theo_khoi"]}

    for r in range(1, len(tbl.rows)):
        khoi_name = tbl.cell(r, 0).text.strip()
        if not khoi_name or khoi_name.lower() in ("cộng", "tổng"): continue

        q4_str = tbl.cell(r, 3).text.strip()
        try: q4_int = int(q4_str)
        except: q4_int = 0
        tong_q4 += q4_int

        # Tìm match trong data
        match = khoi_dict.get(khoi_name.strip().lower())
        if match:
            q1 = match["q1"]
            d  = match["delta"]
            tong_q1 += q1
            set_cell_text_keep_format(tbl.cell(r, 3)._tc, str(q1))
            set_cell_text_keep_format(tbl.cell(r, 4)._tc, delta_str(d))
            # Tô màu delta
            if d is not None and d != 0:
                fg = DELTA_POS_COLOR if d > 0 else DELTA_NEG_COLOR
                tf = tbl.cell(r, 4).text_frame
                for para in tf.paragraphs:
                    for run in para.runs:
                        run.font.color.rgb = fg
                        run.font.bold = True
        else:
            set_cell_text_keep_format(tbl.cell(r, 4)._tc, "—")
            tong_q1 += q4_int

    # Dòng Cộng
    for r in range(1, len(tbl.rows)):
        if tbl.cell(r, 0).text.strip().lower() in ("cộng", "tổng"):
            d_tot = tong_q1 - tong_q4
            set_cell_text_keep_format(tbl.cell(r, 3)._tc, str(tong_q1))
            set_cell_text_keep_format(tbl.cell(r, 4)._tc, delta_str(d_tot))

    # Cập nhật text mô tả
    for shape in s.shapes:
        if shape.has_text_frame and not shape.has_table:
            txt = shape.text_frame.text
            if "505" in txt or "số lượng vbcs" in txt.lower():
                for para in shape.text_frame.paragraphs:
                    for run in para.runs:
                        run.text = run.text \
                            .replace("31/12/2025", f"Q{quy}/{nam}") \
                            .replace("505", str(tong_q1))
    return s

def build_phan1_chart(dest_prs, data, SW, SH):
    """Slide: donut chart phân bổ VB theo Khối"""
    blank = dest_prs.slide_layouts[6]
    s = dest_prs.slides.add_slide(blank)

    add_header(s, "1. Phân bổ VBCS theo Khối", SW=SW)

    chart_data = ChartData()
    khois = data["phan1"]["theo_khoi"]
    chart_data.categories = [r["ten"] for r in khois]
    chart_data.add_series("VBCS", [r["q1"] for r in khois])

    chart_frame = s.shapes.add_chart(
        XL_CHART_TYPE.DOUGHNUT,
        Inches(0.5), Inches(0.7),
        int(SW - Inches(1.0)), int(SH - Inches(1.0)),
        chart_data
    )
    chart = chart_frame.chart
    chart.has_legend = True
    chart.legend.position = XL_LEGEND_POSITION.RIGHT
    chart.legend.include_in_layout = False

    # Tô màu từng slice
    plot = chart.plots[0]
    for i, point in enumerate(plot.series[0].points):
        hex_c = CHART_COLORS[i % len(CHART_COLORS)]
        point.format.fill.solid()
        point.format.fill.fore_color.rgb = rgb(hex_c)

    chart.chart_title.has_text_frame = True
    chart.chart_title.text_frame.text = f"Tổng: {data['phan1']['tong_q1']} VB"

    return s

def build_phan2_bang(src_prs, dest_prs, data, SW, SH):
    """Slide: tổng quan VB>5N theo Khối — clone slide 18"""
    s = clone_slide(src_prs, 17, dest_prs)
    tbl, _ = get_table(s)
    khois    = data["phan2"]["theo_khoi"]
    gn_tot   = sum(r["giu_nguyen"] for r in khois)
    sd_tot   = sum(r["sdbs"] for r in khois)
    th_tot   = sum(r["tuyen_huy"] for r in khois)
    tot_all  = data["phan2"]["tong"]
    quy, nam = data["quy"], data["nam"]

    # Cập nhật text mô tả
    for shape in s.shapes:
        if shape.has_text_frame and not shape.has_table:
            txt = shape.text_frame.text
            if "163" in txt or "05 năm" in txt.lower():
                tf = shape.text_frame
                ns = "http://schemas.openxmlformats.org/drawingml/2006/main"
                # Lấy rPr đầu tiên
                first_rPr = None
                for para in tf.paragraphs:
                    for run in para.runs:
                        rp = run._r.find(f"{{{ns}}}rPr")
                        if rp is not None: first_rPr = copy.deepcopy(rp); break
                    if first_rPr is not None: break
                # Xoá hết, build lại
                while len(tf.paragraphs) > 1:
                    p = tf.paragraphs[-1]._p
                    p.getparent().remove(p)
                p0 = tf.paragraphs[0]._p
                for r in p0.findall(f"{{{ns}}}r"): p0.remove(r)
                lines = [
                    f"Tổng số VBCS đã ban hành > 05 năm: {tot_all} VB, trong đó:",
                    f"{sd_tot}/{tot_all} VB đề xuất SĐBS/thay thế.",
                    f"{th_tot}/{tot_all} VB tuyên hủy/hết hiệu lực.",
                    f"{gn_tot}/{tot_all} VB đề xuất giữ nguyên.",
                ]
                for li, line in enumerate(lines):
                    if li == 0:
                        r_el = etree.SubElement(p0, f"{{{ns}}}r")
                        if first_rPr is not None: r_el.insert(0, copy.deepcopy(first_rPr))
                        t_el = etree.SubElement(r_el, f"{{{ns}}}t"); t_el.text = line
                    else:
                        etree.SubElement(p0, f"{{{ns}}}br")
                        r_el = etree.SubElement(p0, f"{{{ns}}}r")
                        if first_rPr is not None: r_el.insert(0, copy.deepcopy(first_rPr))
                        t_el = etree.SubElement(r_el, f"{{{ns}}}t"); t_el.text = line

    if tbl:
        n_need = len(khois) + 1  # +1 dòng Cộng
        add_rows_to_table(tbl, n_need)

        def v(x): return str(x) if x else "—"
        for i, row in enumerate(khois):
            ri = i + 1
            set_cell_text_keep_format(tbl.cell(ri, 0)._tc, row["ten"])
            set_cell_text_keep_format(tbl.cell(ri, 1)._tc, v(row["giu_nguyen"]))
            set_cell_text_keep_format(tbl.cell(ri, 2)._tc, v(row["sdbs"]))
            set_cell_text_keep_format(tbl.cell(ri, 3)._tc, v(row["tuyen_huy"]))
            set_cell_text_keep_format(tbl.cell(ri, 4)._tc, v(row["tong"]))
            # Tô màu: khối nhiều tuyên hủy → nhạt đỏ
            if row["tuyen_huy"] > 0:
                for c in range(len(tbl.columns)):
                    set_cell_bg(tbl.cell(ri, c)._tc, rgb_tuple("FFE0E0"))
            elif row["sdbs"] > row["giu_nguyen"]:
                for c in range(len(tbl.columns)):
                    set_cell_bg(tbl.cell(ri, c)._tc, rgb_tuple("E4F7FF"))

        # Dòng Cộng
        last = len(tbl.rows) - 1
        set_cell_text_keep_format(tbl.cell(last, 0)._tc, "Cộng")
        set_cell_text_keep_format(tbl.cell(last, 1)._tc, str(gn_tot))
        set_cell_text_keep_format(tbl.cell(last, 2)._tc, str(sd_tot))
        set_cell_text_keep_format(tbl.cell(last, 3)._tc, str(th_tot))
        set_cell_text_keep_format(tbl.cell(last, 4)._tc, str(tot_all))

    return s

def build_phan2_chart(dest_prs, data, SW, SH):
    """Slide: stacked bar chart trạng thái VB>5N theo Khối"""
    blank = dest_prs.slide_layouts[6]
    s = dest_prs.slides.add_slide(blank)
    add_header(s, "2. Trạng thái VBCS trên 5 năm theo Khối", SW=SW)

    khois = data["phan2"]["theo_khoi"]
    chart_data = ChartData()
    chart_data.categories = [r["ten"] for r in khois]
    chart_data.add_series("Giữ nguyên",      [r["giu_nguyen"] for r in khois])
    chart_data.add_series("SĐBS/Thay thế",   [r["sdbs"]       for r in khois])
    chart_data.add_series("Tuyên hủy/HHL",   [r["tuyen_huy"]  for r in khois])

    chart_frame = s.shapes.add_chart(
        XL_CHART_TYPE.BAR_STACKED,
        Inches(0.3), Inches(0.75),
        int(SW - Inches(0.6)), int(SH - Inches(0.95)),
        chart_data
    )
    chart = chart_frame.chart
    chart.has_legend = True
    chart.legend.position = XL_LEGEND_POSITION.BOTTOM

    series_colors = [STYLE["colors"]["blue"], STYLE["colors"]["navy"], STYLE["colors"]["red"]]
    for i, series in enumerate(chart.series):
        series.format.fill.solid()
        series.format.fill.fore_color.rgb = rgb(series_colors[i % len(series_colors)])

    return s

def build_phan2_detail(src_prs, dest_prs, data, SW, SH):
    """Slides: chi tiết VB có tiến độ thực hiện"""
    rows   = data["phan2"]["co_tien_do"]
    slides = []
    PER_PAGE = ROWS_PER_PAGE

    pages = max(1, (len(rows) + PER_PAGE - 1) // PER_PAGE)
    for pi in range(pages):
        chunk = rows[pi*PER_PAGE:(pi+1)*PER_PAGE]
        s = clone_slide(src_prs, 18, dest_prs)  # clone slide 19

        # Cập nhật tiêu đề slide
        for shape in s.shapes:
            if shape.has_text_frame and not shape.has_table:
                txt = shape.text_frame.text
                if "tuyên hủy" in txt.lower() or "trên 5" in txt.lower() or "văn bản" in txt.lower():
                    for para in shape.text_frame.paragraphs:
                        if para.text.strip():
                            for run in para.runs:
                                if run.text.strip():
                                    run.text = f"VB có tiến độ thực hiện — {pi+1}/{pages} (tổng {len(rows)} VB)"
                                    break
                            break

        tbl, _ = get_table(s)
        if tbl:
            # Header
            set_cell_text_keep_format(tbl.cell(0,0)._tc, "TT")
            set_cell_text_keep_format(tbl.cell(0,1)._tc, "Khối")
            set_cell_text_keep_format(tbl.cell(0,2)._tc, "Số ký hiệu VB")
            set_cell_text_keep_format(tbl.cell(0,3)._tc, "Trích yếu")
            set_cell_text_keep_format(tbl.cell(0,4)._tc, "Trạng thái & Tiến độ thực hiện")

            add_rows_to_table(tbl, len(chunk))

            for i, row in enumerate(chunk):
                ri = i + 1
                ts  = row["trang_thai"]
                bg  = TS_COLOR_MAP.get(ts)
                bg_tup = rgb_tuple(bg) if bg else (CT["light_blue"] if i % 2 == 0 else None)

                if bg_tup:
                    for c in range(len(tbl.columns)):
                        set_cell_bg(tbl.cell(ri, c)._tc, bg_tup)

                set_cell_text_keep_format(tbl.cell(ri, 0)._tc, pi*PER_PAGE + i + 1)
                set_cell_text_keep_format(tbl.cell(ri, 1)._tc, row["khoi"])
                set_cell_text_keep_format(tbl.cell(ri, 2)._tc, row["soky"])
                set_cell_text_keep_format(tbl.cell(ri, 3)._tc, row["trich"])

                tien_ghi = row["tien_do"]
                if row["ghi_chu"]:
                    tien_ghi += ("\n— " + row["ghi_chu"]) if tien_ghi else row["ghi_chu"]
                set_cell_text_keep_format(tbl.cell(ri, 4)._tc, f"[{ts}]\n{tien_ghi}")

        slides.append(s)
    return slides


def build_phan3(src_prs, dest_prs, data, SW, SH):
    """Slides: VBCS điều chỉnh theo QĐPL — clone slide 26 (index 25) từ template."""
    phan3 = data.get("phan3")
    if not phan3: return []

    rows = phan3["danh_sach"]
    PER_PAGE = 6
    pages = max(1, (len(rows) + PER_PAGE - 1) // PER_PAGE)
    slides = []

    # Clone slide 26 (index 25) — "Chi tiết các văn bản có tác động đến VBCS"
    try:
        src_slide_idx = 25
        _ = src_prs.slides[src_slide_idx]
    except IndexError:
        src_slide_idx = 18  # fallback

    for pi in range(pages):
        chunk = rows[pi*PER_PAGE:(pi+1)*PER_PAGE]
        s = clone_slide(src_prs, src_slide_idx, dest_prs)

        # Cập nhật tiêu đề
        for shape in s.shapes:
            if shape.has_text_frame and not shape.has_table:
                txt = shape.text_frame.text.strip()
                if "chi tiết" in txt.lower() or "tác động" in txt.lower() or "điều chỉnh" in txt.lower():
                    for para in shape.text_frame.paragraphs:
                        if para.text.strip():
                            for run in para.runs:
                                if run.text.strip():
                                    run.text = f"Chi tiết VBPL tác động đến VBCS — {pi+1}/{pages}"
                                    break
                            break

        tbl, _ = get_table(s)
        if tbl:
            # Header cols: TT | Tên VB | Đơn vị đầu mối | VBCS chịu tác động | Kế hoạch điều chỉnh
            set_cell_text_keep_format(tbl.cell(0, 0)._tc, "TT")
            set_cell_text_keep_format(tbl.cell(0, 1)._tc, "Tên văn bản pháp luật")
            set_cell_text_keep_format(tbl.cell(0, 2)._tc, "Đơn vị đầu mối")
            set_cell_text_keep_format(tbl.cell(0, 3)._tc, "VBCS chịu tác động")
            set_cell_text_keep_format(tbl.cell(0, 4)._tc, "Kế hoạch điều chỉnh VB")

            add_rows_to_table(tbl, len(chunk))

            for i, row in enumerate(chunk):
                ri = i + 1
                bg_tup = CT["light_blue"] if i % 2 == 0 else None
                if bg_tup:
                    for c in range(len(tbl.columns)):
                        set_cell_bg(tbl.cell(ri, c)._tc, bg_tup)

                ten = f"{row['so_vb']}\n{row['ten_vb']}" if row['so_vb'] else row['ten_vb']
                ngay = f"HL: {row['ngay_hl']}" if row['ngay_hl'] else ""
                if ngay: ten += f"\n{ngay}"

                ke_hoach = row["ke_hoach"] or row["note"] or "—"

                set_cell_text_keep_format(tbl.cell(ri, 0)._tc, pi*PER_PAGE + i + 1)
                set_cell_text_keep_format(tbl.cell(ri, 1)._tc, ten)
                set_cell_text_keep_format(tbl.cell(ri, 2)._tc, row["don_vi"])
                set_cell_text_keep_format(tbl.cell(ri, 3)._tc, row["vb_ct"])
                set_cell_text_keep_format(tbl.cell(ri, 4)._tc, ke_hoach)

                # Tô đỏ nhạt nếu có kế hoạch điều chỉnh cụ thể
                if row["ke_hoach"] and len(row["ke_hoach"]) > 5:
                    for c in range(len(tbl.columns)):
                        set_cell_bg(tbl.cell(ri, c)._tc, rgb_tuple("E8F4FF"))

        slides.append(s)
    return slides


# ─────────────────────────────────────────────
# Main
# ─────────────────────────────────────────────

def main():
    if len(sys.argv) < 3:
        print("Usage: python3 ubcs-slide-builder.py <data.json> <template.pptx> [output.pptx]")
        sys.exit(1)

    data_path     = sys.argv[1]
    template_path = sys.argv[2]

    with open(data_path, encoding="utf-8") as f:
        data = json.load(f)

    quy = data["quy"]; nam = data["nam"]
    today = date.today().strftime("%Y-%m-%d")

    if len(sys.argv) > 3:
        out_path = sys.argv[3]
    else:
        out_path = os.path.expanduser(
            f"~/Downloads/Slide hop UBCS Quy {quy} {nam} - {today}.pptx"
        )

    print(f"Build slide Quý {quy}/{nam}...")
    src_prs  = Presentation(template_path)
    dest_prs = Presentation(template_path)
    SW = dest_prs.slide_width
    SH = dest_prs.slide_height

    # Xoá hết slide cũ
    sldIdLst = dest_prs.slides._sldIdLst
    for sld in list(sldIdLst):
        rId = sld.get(qn('r:id'))
        sldIdLst.remove(sld)
        try: dest_prs.part.drop_rel(rId)
        except: pass

    # Build slides
    n = 1
    print(f"  Slide {n}: Tiêu đề"); n += 1
    build_title_slide(src_prs, dest_prs, data)

    print(f"  Slide {n}: Danh mục"); n += 1
    build_muc_luc(src_prs, dest_prs, data, SW, SH)

    print(f"  Slide {n}: Phần 1 — Bìa"); n += 1
    build_section_divider(src_prs, dest_prs, "1", "CÂY VĂN BẢN CHÍNH SÁCH", SW, SH)

    print(f"  Slide {n}: Phần 1 — Bảng cây VBCS"); n += 1
    build_phan1_bang(src_prs, dest_prs, data, SW)

    print(f"  Slide {n}: Phần 1 — Donut chart"); n += 1
    build_phan1_chart(dest_prs, data, SW, SH)

    print(f"  Slide {n}: Phần 2 — Bìa"); n += 1
    build_section_divider(src_prs, dest_prs, "2", "VĂN BẢN CHÍNH SÁCH TRÊN 5 NĂM", SW, SH)

    print(f"  Slide {n}: Phần 2 — Bảng tổng quan VB>5N"); n += 1
    build_phan2_bang(src_prs, dest_prs, data, SW, SH)

    print(f"  Slide {n}: Phần 2 — Stacked bar chart"); n += 1
    build_phan2_chart(dest_prs, data, SW, SH)

    rows2 = data["phan2"]["co_tien_do"]
    pages2 = max(1, (len(rows2) + ROWS_PER_PAGE - 1) // ROWS_PER_PAGE)
    print(f"  Slides {n}–{n+pages2-1}: Phần 2 — Chi tiết tiến độ ({len(rows2)} VB, {pages2} trang)")
    n += pages2
    build_phan2_detail(src_prs, dest_prs, data, SW, SH)

    if data.get("phan3"):
        rows3 = data["phan3"]["danh_sach"]
        pages3 = max(1, (len(rows3) + 5) // 6)
        print(f"  Slide {n}: Phần 3 — Bìa"); n += 1
        build_section_divider(src_prs, dest_prs, "3",
                              "VBCS ĐIỀU CHỈNH THEO QĐPL QUÝ 1/2026", SW, SH)
        print(f"  Slides {n}–{n+pages3-1}: Phần 3 — QĐPL ({len(rows3)} VB, {pages3} trang)")
        build_phan3(src_prs, dest_prs, data, SW, SH)

    dest_prs.save(out_path)
    print(f"\n✓ Saved: {out_path}")
    print(f"  Total slides: {len(dest_prs.slides)}")

if __name__ == "__main__":
    main()
